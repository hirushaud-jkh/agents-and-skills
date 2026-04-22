"""
OCTAVE Deck Generator — Core Helpers
=====================================
Load at the top of every deck generation script:
    exec(open(os.path.expanduser("~/.claude/skills/deck/recipes.py")).read())

Use these building blocks freely to compose any slide layout.
Combine shapes, boxes, arrows, circles, tables creatively — no fixed templates.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# =============================================
# COLOR PALETTE
# =============================================
G = RGBColor(0x26, 0xEA, 0xA0)    # Brand green
DG = RGBColor(0x1A, 0xA5, 0x70)   # Darker green
W = RGBColor(0xFF, 0xFF, 0xFF)    # White
BK = RGBColor(0x2D, 0x2D, 0x2D)   # Near black
MG = RGBColor(0x4A, 0x4A, 0x4A)   # Medium gray
LG = RGBColor(0xF0, 0xF0, 0xF0)   # Light gray
BLU = RGBColor(0x3A, 0x86, 0xFF)  # Blue accent
ORG = RGBColor(0xFF, 0x8C, 0x42)  # Orange accent
PRP = RGBColor(0x9B, 0x5D, 0xE5)  # Purple accent
RED = RGBColor(0xE8, 0x4D, 0x4D)  # Red accent
TEAL = RGBColor(0x20, 0xC9, 0x97) # Teal
DRED = RGBColor(0xC0, 0x39, 0x39) # Dark red


# =============================================
# BRAND FONT
# =============================================
FONT = "Lato"

# =============================================
# CORE HELPERS
# =============================================

def sf(run, size=None, bold=None, color=None, italic=None):
    """Style a font run — always sets brand font."""
    run.font.name = FONT
    if size: run.font.size = Pt(size)
    if bold is not None: run.font.bold = bold
    if color: run.font.color.rgb = color
    if italic is not None: run.font.italic = italic


def box(sl, l, t, w, h, text, fill=G, fc=W, fs=12, bold=True, align=PP_ALIGN.CENTER, radius=0.12):
    """Rounded rectangle box with text."""
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    if radius: s.adjustments[0] = radius
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; sf(r, fs, bold, fc)
    return s


def obox(sl, l, t, w, h, text, border_color=G, fc=BK, fs=12, bold=False, align=PP_ALIGN.CENTER, radius=0.08, border_width=1.5):
    """Outlined box (no fill, colored border) — for cards, callouts, feature boxes."""
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = W
    s.line.color.rgb = border_color; s.line.width = Pt(border_width)
    if radius: s.adjustments[0] = radius
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; sf(r, fs, bold, fc)
    return s


def mbox(sl, l, t, w, h, lines, fill=G, fc=W, fs=11, align=PP_ALIGN.CENTER, radius=0.06):
    """Multi-line box. lines = [(text, size, bold, color), ...]"""
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    if radius: s.adjustments[0] = radius
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    for i, (txt_val, sz, b, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; r = p.add_run(); r.text = txt_val; sf(r, sz, b, c)
    return s


def txt(sl, l, t, w, h, text, fs=12, color=MG, bold=False, align=PP_ALIGN.LEFT, italic=False):
    """Text box."""
    s = sl.shapes.add_textbox(l, t, w, h)
    tf = s.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; sf(r, fs, bold, color, italic)
    return s


def mtxt(sl, l, t, w, h, lines, align=PP_ALIGN.LEFT):
    """Multi-line text box. lines = [(text, size, bold, color), ...]"""
    s = sl.shapes.add_textbox(l, t, w, h)
    tf = s.text_frame; tf.word_wrap = True
    for i, (txt_val, sz, b, c) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; r = p.add_run(); r.text = txt_val; sf(r, sz, b, c)
    return s


def ar(sl, l, t, w=Inches(0.3), h=Inches(0.18)):
    """Right arrow."""
    s = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = G; s.line.fill.background()
    return s


def ad(sl, l, t, w=Inches(0.2), h=Inches(0.25)):
    """Down arrow."""
    s = sl.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = G; s.line.fill.background()
    return s


def circ(sl, l, t, sz, fill=G, text="", fc=W, fs=12):
    """Circle with optional text."""
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    if text:
        tf = s.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text; sf(r, fs, True, fc)
    return s


def icon_box(sl, l, t, w, h, icon_text, title, desc, icon_color=G, icon_size=Inches(0.5)):
    """Card with a colored icon circle, title, and description — good for feature/benefit cards."""
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = LG; s.line.fill.background()
    s.adjustments[0] = 0.06
    circ(sl, l + Inches(0.15), t + Inches(0.15), icon_size, icon_color, icon_text, W, 14)
    txt(sl, l + icon_size + Inches(0.3), t + Inches(0.15), w - icon_size - Inches(0.45), Inches(0.35),
        title, 13, BK, True)
    txt(sl, l + Inches(0.15), t + icon_size + Inches(0.3), w - Inches(0.3), h - icon_size - Inches(0.45),
        desc, 11, MG)
    return s


def accent_bar(sl, l, t, w, color=G, thickness=Inches(0.06)):
    """Thin colored accent bar — use at top of cards for visual hierarchy."""
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, thickness)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def hline(sl, l, t, w, color=G):
    """Horizontal line."""
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Inches(0.025))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s


def badge(sl, l, t, text, color=G, fc=W, fs=10, padding=0.15):
    """Small pill-shaped badge/tag — for labels, status indicators, categories."""
    w = Inches(len(text) * 0.09 + padding * 2)
    h = Inches(0.32)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    s.adjustments[0] = 0.5  # fully rounded ends
    tf = s.text_frame; tf.word_wrap = False
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; sf(r, fs, True, fc)
    return s


def connector(sl, l1, t1, l2, t2, color=G, width=1.5):
    """Simple line connector between two points."""
    s = sl.shapes.add_connector(1, l1, t1, l2, t2)  # 1 = straight
    s.line.color.rgb = color; s.line.width = Pt(width)
    return s


def stbl(sl, data, l, t, w, h, cw=None):
    """Styled table with brand header. Headers: 13pt bold white on green, Body: 11pt, alternating rows."""
    rows, cols = len(data), len(data[0])
    ts = sl.shapes.add_table(rows, cols, l, t, w, h); tbl = ts.table
    if cw:
        for i, c in enumerate(cw): tbl.columns[i].width = c
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = ""
            cell.vertical_anchor = 1
            p = cell.text_frame.paragraphs[0]; r = p.add_run(); r.text = str(val)
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = G
                sf(r, 13, True, W)
                cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
                cell.margin_top = Inches(0.08); cell.margin_bottom = Inches(0.08)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = W if ri % 2 == 1 else LG
                sf(r, 11, False, BK)
                cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.12)
                cell.margin_top = Inches(0.06); cell.margin_bottom = Inches(0.06)
            p.alignment = PP_ALIGN.LEFT
    return ts


def bullets(sl, l, t, w, h, items, fs=12, color=BK, bullet_color=G, spacing=Pt(6)):
    """Bulleted list. items = list of strings or [(text, size, bold, color), ...] tuples."""
    s = sl.shapes.add_textbox(l, t, w, h)
    tf = s.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = spacing
        br = p.add_run(); br.text = "\u2022  "; sf(br, fs, True, bullet_color)
        if isinstance(item, str):
            r = p.add_run(); r.text = item; sf(r, fs, False, color)
        else:
            txt_val, sz, b, c = item
            r = p.add_run(); r.text = txt_val; sf(r, sz, b, c)
    return s


def img(sl, path, l, t, w=None, h=None):
    """Add image. Provide w or h (or both) — aspect ratio preserved if only one given."""
    return sl.shapes.add_picture(path, l, t, w, h)


def new_deck(template_path):
    """Load template and remove all sample slides — returns clean Presentation with only layouts.
    Use instead of Presentation(template_path) when generating a new deck."""
    prs = Presentation(template_path)
    ns_r = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'
    ns_p = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    sldIdLst = prs.part._element.find(ns_p + 'sldIdLst')
    for sldId in list(sldIdLst):
        rId = sldId.get(ns_r + 'id')
        sldIdLst.remove(sldId)
        prs.part.drop_rel(rId)
    return prs


def set_title(sl, title_text):
    """Set slide title and REMOVE unused placeholders (not just clear them)."""
    spTree = sl.shapes._spTree
    remove = []
    for ph in sl.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
            r = ph.text_frame.paragraphs[0].add_run()
            r.text = title_text; sf(r, 28, True, BK)
        elif ph.placeholder_format.idx in (1, 10):
            remove.append(ph._element)
    for el in remove:
        spTree.remove(el)
